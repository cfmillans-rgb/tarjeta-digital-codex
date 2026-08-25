import collections
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Simple styling
    title_shape.text_frame.paragraphs[0].font.bold = True
    return slide

def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    tf = body_shape.text_frame
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            
    return slide

def add_two_content_slide(prs, title, col1, col2):
    slide_layout = prs.slide_layouts[3] # Two Content
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape1 = slide.placeholders[1]
    body_shape2 = slide.placeholders[2]
    
    title_shape.text = title
    
    tf1 = body_shape1.text_frame
    for i, item in enumerate(col1):
        if i == 0:
            p = tf1.paragraphs[0]
            p.text = item
        else:
            p = tf1.add_paragraph()
            p.text = item
            p.level = 0
            
    tf2 = body_shape2.text_frame
    for i, item in enumerate(col2):
        if i == 0:
            p = tf2.paragraphs[0]
            p.text = item
        else:
            p = tf2.add_paragraph()
            p.text = item
            p.level = 0
            
    return slide

def main():
    # Attempt to monkey patch collections for pptx if needed on python 3.10+
    try:
        setattr(collections, 'Container', collections.abc.Container)
        setattr(collections, 'Mapping', collections.abc.Mapping)
        setattr(collections, 'MutableMapping', collections.abc.MutableMapping)
        setattr(collections, 'Sequence', collections.abc.Sequence)
        setattr(collections, 'Iterable', collections.abc.Iterable)
        setattr(collections, 'Iterator', collections.abc.Iterator)
    except AttributeError:
        pass
        
    prs = Presentation()
    
    # Slide 1: Title
    add_title_slide(prs, "Digitaliza Lab", "Capital Abeja Emprende 2026\nManual de Presentación y Defensa")
    
    # Slide 2: Mi Empresa
    add_bullet_slide(prs, "1. Mi Empresa", [
        "Rubro: Servicios informáticos.",
        "¿Qué hacemos? Desarrollo de landing pages profesionales orientadas a la conversión.",
        "Objetivo: Ayudar a profesionales, emprendedores y microempresas a conseguir más clientes mediante una presencia digital rápida y efectiva."
    ])
    
    # Slide 3: El Problema
    add_bullet_slide(prs, "El Problema que Resolvemos", [
        "Muchas microempresas no tienen página web porque:",
        "• Las agencias tradicionales son costosas.",
        "• Los tiempos de entrega son largos.",
        "• Crear una página por cuenta propia requiere conocimientos técnicos.",
        "Digitaliza Lab elimina esas barreras ofreciendo una solución rápida, simple y accesible."
    ])
    
    # Slide 4: Propuesta de Valor
    add_bullet_slide(prs, "2. Propuesta de Valor", [
        "Landing page profesional.",
        "Entrega en 72 horas hábiles.",
        "Precio fijo de $99.990.",
        "Todo incluido en un solo servicio (sin costos ocultos):",
        "• Dominio incluido.",
        "• Hosting incluido.",
        "• Certificado SSL.",
        "• Correo corporativo por un año."
    ])
    
    # Slide 5: Clientes
    add_two_content_slide(prs, "3. Nuestros Clientes", [
        "Profesionales independientes",
        "• Abogados, psicólogos, arquitectos, etc.",
        "• Necesitan mejorar su imagen profesional.",
        "",
        "Microempresas",
        "• Construcción, transporte, servicios técnicos, etc.",
        "• Necesitan captar clientes."
    ], [
        "Emprendedores",
        "• Velas, cosmética, tiendas.",
        "• Necesitan vender por internet.",
        "",
        "Punto en común:",
        "No poseen una presencia digital profesional."
    ])
    
    # Slide 6: Diferenciadores
    add_bullet_slide(prs, "4. Diferenciadores", [
        "1. Rapidez: Entrega garantizada en 72 horas.",
        "2. Todo incluido: No existen cobros separados. El cliente recibe dominio, hosting, SSL y correo.",
        "3. Conversión: No se diseña solamente una página bonita. Se construye una herramienta para conseguir clientes."
    ])
    
    # Slide 7: Validación y Modelo de Negocio
    add_two_content_slide(prs, "5. y 6. Validación y Modelo", [
        "Validación",
        "• El negocio ya fue probado en distintos rubros.",
        "• Se validaron tiempos, proceso, calidad y escalabilidad.",
        "• El subsidio es para fortalecer y escalar un negocio que ya funciona."
    ], [
        "Modelo de Negocio",
        "• Venta principal: Landing Page a $99.990",
        "• Ingresos futuros:",
        "  - Mantenciones",
        "  - Renovaciones de dominio y hosting",
        "  - Correo corporativo adicional"
    ])
    
    # Slide 8: Operación
    add_two_content_slide(prs, "7., 8. y 9. Operación", [
        "Proceso de Trabajo",
        "1. Cliente solicita vía Formulario",
        "2. Diseño y Programación",
        "3. Revisión con el cliente",
        "4. Publicación y Entrega",
        "5. Soporte y Postventa"
    ], [
        "Canales y Alianzas",
        "• Canales: Sitio web, Instagram, Facebook, Mercado Libre, WhatsApp Business.",
        "• Alianzas clave:",
        "  - NIC Chile y Proveedor Hosting.",
        "  - Mercado Pago.",
        "  - Profesionales (Contadores, Fotógrafos, Diseñadores)."
    ])
    
    # Slide 9: Sustentabilidad
    add_two_content_slide(prs, "10. Sustentabilidad", [
        "Impacto Actual",
        "• Servicio 100% digital.",
        "• Cero uso de papel.",
        "• Sin traslados físicos.",
        "• Menor huella de carbono."
    ], [
        "Proyección a Futuro",
        "• Migración a Green Hosting.",
        "• Uso de equipos más eficientes energéticamente.",
        "• Fomento del reciclaje electrónico."
    ])
    
    # Slide 10: Presupuesto
    add_bullet_slide(prs, "11. Presupuesto (Total: $3.500.000)", [
        "• Gestión empresarial ($500.000): Constitución empresa, Meta Ads + Google Ads.",
        "• Activos Fijos ($1.760.000): Notebook profesional, Monitor 27 pulgadas, Mouse y Teclado ergonómicos.",
        "• Activos Intangibles ($740.000): ChatGPT Plus, Cursor Pro, Canva Pro, Google Workspace, Registro INAPI.",
        "• Capital de Trabajo ($500.000): Dominios, Hosting, Envato Elements.",
        "",
        "Nota: Las herramientas (Cursor Pro, Canva, Workspace) fortalecen directamente la productividad del negocio."
    ])
    
    # Slide 11: Justificación de Inversión
    add_bullet_slide(prs, "12. ¿Por qué cada compra?", [
        "• Notebook y Monitor: Estación principal de trabajo para desarrollo, IA y diseño concurrente con mayor espacio visual.",
        "• Ergonomía (Mouse/Teclado): Prevención de lesiones en largas jornadas de programación.",
        "• IA y Herramientas (ChatGPT Plus, Cursor Pro): Multiplica la productividad en código, SEO y textos comerciales.",
        "• Canva Pro y Workspace: Diseño ágil para redes y gestión corporativa formal.",
        "• Registro Marca INAPI: Protección y formalidad legal."
    ])
    
    # Slide 12: Objetivo y Frase de Cierre
    add_bullet_slide(prs, "Objetivo del Proyecto y Conclusión", [
        "Objetivo del Proyecto:",
        "Formalizar la empresa, fortalecer la infraestructura, profesionalizar la operación y escalar ventas mediante publicidad.",
        "",
        "Conclusión:",
        '"El subsidio no está destinado a crear la idea de negocio desde cero, sino a fortalecer una empresa que ya validó su servicio. La inversión permitirá formalizar Digitaliza Lab, mejorar su capacidad operativa, aumentar la productividad y captar más clientes para consolidar un negocio sostenible en el tiempo."'
    ])
    
    prs.save("Digitaliza_Lab_Pitch.pptx")
    print("Presentation saved successfully as Digitaliza_Lab_Pitch.pptx")

if __name__ == "__main__":
    main()
